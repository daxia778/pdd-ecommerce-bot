"""
NotebookLM 客户端 — 统一调用接口。

调用层级：
  task_coordinator → NotebookLMClient.generate_ppt()
                   → NotebookLMPlaywrightAutomation（真实浏览器自动化）
                   → 降级 → python-pptx 本地生成

generate_ppt() 返回值：可直接通过 HTTP 访问的文件 URL（相对于本服务 /files/ 路径）。
"""

import asyncio
import os

from src.services.notebooklm_playwright import NotebookLMPlaywrightAutomation
from src.utils.logger import logger

# 本服务对外提供文件下载的 URL 前缀（对应 main.py 中挂载的 /files 静态目录）
_FILES_URL_PREFIX = os.environ.get("PDD_BOT_BASE_URL", "http://localhost:8100") + "/files"

# 输出目录（与 playwright 服务保持一致）
_OUTPUT_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "data", "output"))


class NotebookLMClient:
    def __init__(self):
        self.playwright_service = NotebookLMPlaywrightAutomation()

    async def generate_ppt(
        self,
        topic: str,
        pages: int,
        style: str,
        requirements: str = "",
    ) -> str:
        """
        生成 PPT 并返回可供下载的 HTTP URL。

        优先用 Playwright 自动化 NotebookLM（需已登录），
        失败时降级为 python-pptx 本地生成草稿版。
        """
        logger.info(f"NotebookLM | 开始生成 | 主题: {topic} | 页数: {pages} | 风格: {style}")

        # ── 方案 A: Playwright 自动化 ─────────────────────────────────────────
        try:
            local_path = await self.playwright_service.generate_ppt(
                topic=topic,
                pages=pages,
                style=style,
                requirements=requirements,
            )
            file_url = self._local_path_to_url(local_path)
            logger.info(f"NotebookLM | ✅ Playwright 生成成功: {file_url}")
            return file_url

        except Exception as e:
            logger.warning(f"NotebookLM | Playwright 生成失败，降级到 python-pptx: {e}")

        # ── 方案 B: python-pptx 本地降级草稿 ────────────────────────────────
        return await self._fallback_pptx(topic, pages, requirements)

    @staticmethod
    def _local_path_to_url(local_path: str) -> str:
        """将本地绝对路径转换为可供外部访问的 /files/ URL"""
        filename = os.path.basename(local_path)
        return f"{_FILES_URL_PREFIX}/{filename}"

    async def _fallback_pptx(self, topic: str, pages: int, requirements: str) -> str:
        """
        使用 python-pptx 生成简单草稿 PPT（不依赖网络/浏览器）。
        主要作为降级保底，确保流水线不会因 Playwright 问题完全阻塞。
        """
        try:
            import time

            from pptx import Presentation

            prs = Presentation()

            # 封面页
            cover_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(cover_layout)
            slide.shapes.title.text = topic
            if slide.placeholders and len(slide.placeholders) > 1:
                slide.placeholders[1].text = f"需求说明：{requirements}\n页数：{pages} 页（草稿版）"

            # 正文占位页
            content_layout = prs.slide_layouts[1]
            for i in range(1, min(pages, 10)):  # 最多生成 10 页
                s = prs.slides.add_slide(content_layout)
                s.shapes.title.text = f"第 {i + 1} 页 — {topic}"
                if s.placeholders and len(s.placeholders) > 1:
                    s.placeholders[1].text = f"• 本页为草稿占位内容\n• 请由设计师根据主题补充实际内容\n• 主题：{topic}"

            os.makedirs(_OUTPUT_DIR, exist_ok=True)
            safe_name = "".join(c for c in topic[:20] if c.isalnum() or c in " _").replace(" ", "_")
            filename = f"draft_{safe_name}_{int(time.time())}.pptx"
            file_path = os.path.join(_OUTPUT_DIR, filename)
            prs.save(file_path)

            file_url = NotebookLMClient._local_path_to_url(file_path)
            logger.info(f"NotebookLM | python-pptx 草稿已生成: {file_url}")
            return file_url

        except ImportError:
            logger.error("python-pptx 未安装，请执行: pip install python-pptx")
            raise
        except Exception as e:
            logger.error(f"NotebookLM | python-pptx 降级也失败: {e}")
            raise

    async def remove_watermark(self, file_url: str) -> str:
        """
        后处理：去除水印（当前为透传，预留给未来集成水印工具）。
        NotebookLM 直接下载的 PPTX 无水印，此方法直接返回原 URL。
        """
        logger.info(f"WatermarkRemover | 文件 URL: {file_url}（NotebookLM 原版无水印，直接透传）")
        await asyncio.sleep(0.1)
        return file_url


# 模块级单例
notebooklm_client = NotebookLMClient()
